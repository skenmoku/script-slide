import io
import os
import re
import tempfile
from typing import List, Tuple

from flask import (
    Flask, request, render_template, send_file, abort, jsonify
)
from werkzeug.utils import secure_filename
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# --------------------------------------------------
# Flask 基本設定（本番想定）
# --------------------------------------------------
app = Flask(__name__, template_folder="templates")
app.config["MAX_CONTENT_LENGTH"] = 30 * 1024 * 1024  # 30MB
ALLOWED_EXT = {".pptx"}

# --------------------------------------------------
# レイアウト / 表示 設定
# --------------------------------------------------

# 1スライドあたりの最大文字数（チャンク境界）
MAX_CHARS_PER_SLIDE = 150

# スライドサイズ (16:9)
SLIDE_WIDTH_CM = 33.867
SLIDE_HEIGHT_CM = 19.05

# ★ご指定のテキストボックスサイズ
TEXT_LEFT_CM = 0.80
TEXT_TOP_CM = 0.80
TEXT_WIDTH_CM = 32.27   # 指定：32.27cm
TEXT_HEIGHT_CM = 17.32  # 指定：17.32cm

# 右下の枠（フォト置き場イメージ）
FRAME_LEFT_CM = 25.87
FRAME_TOP_CM = 14.55
FRAME_WIDTH_CM = 8.0
FRAME_HEIGHT_CM = 4.5

# 分割時のみ表示するページインジケータ
PAGE_LEFT_CM = 21.94
PAGE_TOP_CM = 16.93
PAGE_COLOR = RGBColor(0x00, 0x9D, 0xFF)  # #009DFF
PAGE_FONT_SIZE_PT = 32
PAGE_FONT_BOLD = True

# 話者ごとの明示色（固定）
NAME_FIXED_COLORS = {
    "仲條": RGBColor(0x00, 0xFD, 0xFF),  # 水色
    "三村": RGBColor(0xFF, 0xFF, 0xFF),  # 白
    "星野": RGBColor(0xFF, 0xFF, 0x00),  # 黄
}

# 自動割当用の色プール（循環）
AUTO_COLOR_POOL = [
    RGBColor(0xFF, 0x40, 0xFF),  # ピンク
    RGBColor(0xFF, 0xA5, 0x00),  # オレンジ
    RGBColor(0xFF, 0xFB, 0x00),  # 黄（サブ）
]
name_color_map = {}
_auto_color_idx = 0


# --------------------------------------------------
# ユーティリティ
# --------------------------------------------------
def _ext_ok(filename: str) -> bool:
    _, ext = os.path.splitext(filename.lower())
    return ext in ALLOWED_EXT

def clean_text(text: str) -> str:
    """ノート欄の制御文字を除去"""
    return text.replace("\x0b", "").strip()

def get_color_for_name(name: str) -> RGBColor:
    """話者名に応じた色を返す（固定→既存→自動）"""
    global _auto_color_idx
    if not name:
        return RGBColor(0xFF, 0xFF, 0xFF)  # デフォルト白

    if name in NAME_FIXED_COLORS:
        return NAME_FIXED_COLORS[name]
    if name in name_color_map:
        return name_color_map[name]

    color = AUTO_COLOR_POOL[_auto_color_idx % len(AUTO_COLOR_POOL)]
    name_color_map[name] = color
    _auto_color_idx += 1
    return color

def parse_notes_into_segments(note_text: str) -> List[Tuple[str, str]]:
    """
    ノートを《名前》単位で抽出 -> [(name, text)]。
    連続する同一話者は結合。
    """
    segments = []
    current_name = None
    buffer = []

    for raw in note_text.splitlines():
        line = raw.strip()
        if not line:
            continue

        m = re.match(r"《(.+?)》", line)
        if m:
            if buffer:
                joined = "".join(buffer).strip()
                if joined:
                    segments.append((current_name, joined))
                buffer = []
            current_name = m.group(1)
            rest = line[m.end():]
            if rest:
                buffer.append(rest)
        else:
            buffer.append(line)

    if buffer:
        joined = "".join(buffer).strip()
        if joined:
            segments.append((current_name, joined))

    # 隣接する同名を統合
    merged = []
    for name, text in segments:
        if merged and merged[-1][0] == name:
            merged[-1] = (name, merged[-1][1] + text)
        else:
            merged.append((name, text))
    return merged

def _split_preserving_words(s: str, limit: int):
    """
    文字列 s を limit 長以下の塊に分割（可能なら単語境界で）。
    日本語主体でも、句読点や空白で切れるときは切る。
    """
    if len(s) <= limit:
        return [s]

    parts = []
    start = 0
    while start < len(s):
        end = min(len(s), start + limit)
        # 余裕がある場合は近くの区切り（空白・句読点）まで戻す
        cut = end
        window = s[start:end]
        m = re.search(r"[、。．，,.?!\s]+(?=[^、。．，,.?!\s]*$)", window)
        if m:
            cut = start + m.end()
        parts.append(s[start:cut])
        start = cut
    return parts

def pack_segments_into_chunks(segments: List[Tuple[str, str]], max_len=MAX_CHARS_PER_SLIDE):
    """
    同じスライド内で話者を分断しない方針は維持しつつ、
    必要ならその話者テキストを単語境界優先で分割。
    """
    chunks: List[List[Tuple[str, str]]] = []
    cur: List[Tuple[str, str]] = []
    cur_len = 0

    def flush():
        nonlocal cur, cur_len
        if cur:
            chunks.append(cur)
            cur = []
            cur_len = 0

    for name, text in segments:
        # その話者のテキストを分割して順に積む
        pieces = _split_preserving_words(text, max_len)
        for piece in pieces:
            piece = piece.strip()
            if not piece:
                continue
            need = len(piece)
            remain = max_len - cur_len
            if need > remain and cur:
                flush()
            cur.append((name, piece))
            cur_len += len(piece)
            if cur_len >= max_len:
                flush()

    flush()
    return chunks

def add_photo_frame(slide):
    """右下に固定枠"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Cm(FRAME_LEFT_CM), Cm(FRAME_TOP_CM), Cm(FRAME_WIDTH_CM), Cm(FRAME_HEIGHT_CM)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xF0)
    shape.line.color.rgb = RGBColor(0x64, 0x64, 0x64)
    shape.line.width = Pt(2)
    return shape

def add_page_indicator(slide, index: int, total: int):
    """分割時のみページ番号を表示"""
    if total <= 1:
        return
    tb = slide.shapes.add_textbox(Cm(PAGE_LEFT_CM), Cm(PAGE_TOP_CM), Cm(4), Cm(1.5))
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = f"{index}/{total}"
    p.alignment = PP_ALIGN.LEFT
    font = p.font
    font.name = "メイリオ"  # 無ければ OS 既定フォントにフォールバック
    font.size = Pt(PAGE_FONT_SIZE_PT)
    font.bold = PAGE_FONT_BOLD
    font.color.rgb = PAGE_COLOR

def create_script_slides(notes: List[str]) -> Presentation:
    prs = Presentation()
    prs.slide_width = Cm(SLIDE_WIDTH_CM)
    prs.slide_height = Cm(SLIDE_HEIGHT_CM)

    for note in notes:
        segments = parse_notes_into_segments(note)
        chunks = pack_segments_into_chunks(segments, MAX_CHARS_PER_SLIDE)
        total_parts = len(chunks)

        for idx, chunk in enumerate(chunks, start=1):
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            # 背景：黒
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(0, 0, 0)

            # テキストボックス（指定サイズ）
            tx = slide.shapes.add_textbox(
                Cm(TEXT_LEFT_CM), Cm(TEXT_TOP_CM), Cm(TEXT_WIDTH_CM), Cm(TEXT_HEIGHT_CM)
            )
            tf = tx.text_frame
            tf.clear()
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.space_before = Pt(0)
            p.space_after = Pt(0)

            # 各話者 run を追加
            for name, part in chunk:
                run = p.add_run()
                prefix = f"《{name}》" if name else ""
                run.text = prefix + part
                f = run.font
                f.name = "メイリオ"
                f.size = Pt(40)
                f.bold = True
                f.color.rgb = get_color_for_name(name)

            add_photo_frame(slide)
            add_page_indicator(slide, idx, total_parts)

    return prs


# --------------------------------------------------
# ルーティング
# --------------------------------------------------
@app.get("/healthz")
def healthz():
    return jsonify(ok=True)

@app.route("/", methods=["GET", "POST"])
def index():
    if request.method == "GET":
        return render_template("index.html")

    # POST: ファイル受領
    if "pptx_file" not in request.files:
        abort(400, description="ファイルが送信されていません。")
    up = request.files["pptx_file"]
    if not up or up.filename == "":
        abort(400, description="ファイル名が空です。")

    filename = secure_filename(up.filename)
    if not _ext_ok(filename):
        abort(400, description="拡張子が .pptx のファイルのみ対応しています。")

    with tempfile.TemporaryDirectory() as td:
        path_in = os.path.join(td, "input.pptx")
        up.save(path_in)

        try:
            prs_in = Presentation(path_in)
        except Exception:
            abort(400, description="PPTX の読み込みに失敗しました。ファイルを確認してください。")

        notes: List[str] = []
        for s in prs_in.slides:
            if s.has_notes_slide and s.notes_slide.notes_text_frame:
                text = s.notes_slide.notes_text_frame.text or ""
                cleaned = clean_text(text)
                if cleaned:
                    notes.append(cleaned)

        # ノートが空の場合も空のプレゼンを返す（1ページも生成されない）
        prs_out = create_script_slides(notes)
        buf = io.BytesIO()
        prs_out.save(buf)
        buf.seek(0)
        return send_file(
            buf,
            as_attachment=True,
            download_name="スクリプトスライド_自動生成.pptx",
            mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )

# ローカル実行用（Render では gunicorn が使われる）
if __name__ == "__main__":
    port = int(os.environ.get("PORT", "5000"))
    app.run(host="0.0.0.0", port=port, debug=False)
